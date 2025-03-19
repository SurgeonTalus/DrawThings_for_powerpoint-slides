import os
import json
import base64
import time
import subprocess
import requests
from pptx import Presentation
from pptx.util import Cm
import glob
import tkinter as tk
from tkinter import filedialog
from pptx.util import Inches

# API Endpoints
LM_STUDIO_URL = "http://localhost:1234/v1/chat/completions"
DRAW_THINGS_URL = "http://127.0.0.1:7860/sdapi/v1/img2img"

# Image generation settings
STEPS = 4  # Set to 4 as requested
REMOVE_BACKGROUND_SHORTCUT = "RemoveBackground"

# Function to extract text from slides
def extract_text_from_pptx(pptx_path):
    presentation = Presentation(pptx_path)
    slides_text = []
    
    for slide in presentation.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text.strip())
        slides_text.append("\n".join(slide_text))

        # Define the subtitleForTTS variable as the slide text
        subtitleForTTS = "\n".join(slide_text)  # Combine all the text on the slide into one string

        # Define the file paths in the Downloads folder
        downloads_folder = os.path.expanduser("~/Downloads")
        text_file_path = os.path.join(downloads_folder, 'plain_text.txt')
        audio_file_path = os.path.join(downloads_folder, 'audio_output.aiff')

        # Save the plain text to a .txt file
        with open(text_file_path, 'w', encoding='utf-8') as text_file:
            text_file.write(subtitleForTTS)

        # Use macOS 'say' command to generate an audio file from the text file
        os.system(f"say -v Henrik -f {text_file_path} -o {audio_file_path}")

        # Add the audio file (embedding audio)
        audio_shape = slide.shapes.add_movie(
            audio_file_path,
            left=Inches(4.17),
            top=Inches(6.74),
            width=Inches(1.67),
            height=Inches(0.76),
            poster_frame_image=None,
            mime_type='audio/aiff'  # Adjust mime type based on the audio file format
        )
    
    
    
    return presentation, slides_text

# Function to get an image description from LM Studio
def description_prompt_text(prompt_text):
    request_data = {
        "messages": [
            {"role": "system", "content": "Do not give chat answers. Only aswer with the description. You are a prompt engenere for stabile diffusion SDXL generating a single prompt. Context bleading is verry comon. Therefore, avoid bad prompts and do not add unessesary complexity. Using the context of the PowerPoint slide. Illustrate examples from the context. All imges have colors."},
            {"role": "user", "content": prompt_text}
        ],
        "model": "gemma-3-12b-it",
        "temperature": 0.1,
        "max_tokens": 500,
        "stream": False,
        "reset_history": True,
    }

    response = requests.post(LM_STUDIO_URL, json=request_data)
    
    if response.status_code == 200:
        response_json = response.json()
        return response_json.get("choices", [{}])[0].get("message", {}).get("content", "No description generated.")
    else:
        print(f"Error generating description. Status code: {response.status_code}")
        return None

# Function to generate an image using Draw Things API
def generate_image(prompt):
    print(f"Generating image with {STEPS} steps...")
    
    params = {
        "prompt": prompt,
        "negative_prompt": "(bokeh, worst quality, low quality, normal quality, (variations):1.4), blur:1.5",
        "seed": 4068245935,
        "steps": STEPS,
        "guidance_scale": 10,
        "batch_count": 1
    }

    headers = {"Content-Type": "application/json"}
    response = requests.post(DRAW_THINGS_URL, json=params, headers=headers)

    if response.status_code == 200:
        data = response.json()
        images = data.get("images", [])
        if images:
            temp_image_path = os.path.join("/tmp", "generated_image.png")
            with open(temp_image_path, "wb") as img_file:
                img_file.write(base64.b64decode(images[0]))

            return temp_image_path

    print(f"Error generating image: {response.status_code}, {response.text}")
    return None

# Function to copy image to clipboard on macOS
def copy_image_to_clipboard(image_path):
    try:
        subprocess.run(["osascript", "-e", f'set the clipboard to (read (POSIX file "{image_path}") as JPEG picture)'], check=True)
        print("Image copied to clipboard.")
    except subprocess.CalledProcessError as e:
        print(f"Error copying image to clipboard: {e}")

# Function to run Siri Shortcut
def run_siri_shortcut(shortcut_name):
    try:
        subprocess.run(["shortcuts", "run", shortcut_name], check=True)
        print(f"Siri Shortcut '{shortcut_name}' executed.")
    except subprocess.CalledProcessError as e:
        print(f"Error running Siri Shortcut: {e}")

# Function to get the latest image from the Downloads folder
def get_latest_downloaded_image():
    list_of_files = glob.glob(os.path.join(os.path.expanduser("~"), "Downloads", "*.png"))  # Get all PNG files
    if not list_of_files:
        return None
    latest_file = max(list_of_files, key=os.path.getctime)  # Get the most recently modified file
    return latest_file

def insert_image_to_slide(slide, image_path, x, y):
    if not os.path.exists(image_path):
        print(f"Image not found: {image_path}")
        return
    
    image_width = Cm(19.05)
    image_height = Cm(19.05)
    slide.shapes.add_picture(image_path, x, y, width=image_width, height=image_height)

# Configuration options
addOriginal = False  # Set to False to skip adding the original image
addSeparated = True  # Set to False to skip adding the background-removed image

# Function to select the PowerPoint file using tkinter
def select_pptx_file():
    root = tk.Tk()
    root.withdraw()  # Hide the main tkinter window
    pptx_file = filedialog.askopenfilename(
        title="Select a PowerPoint File",
        filetypes=[("PowerPoint Files", "*.pptx")]
    )
    return pptx_file

# Main script execution
if __name__ == "__main__":
    PPTX_FILE = select_pptx_file()  # Use tkinter to select the PowerPoint file
    if not PPTX_FILE:
        print("No file selected. Exiting...")
        exit(1)

    presentation, slides = extract_text_from_pptx(PPTX_FILE)

    # Use the selected file's name to create a new filename for saving the output
    base_name = os.path.splitext(os.path.basename(PPTX_FILE))[0]  # Extract file name without extension
    output_pptx_file = os.path.join(os.path.dirname(PPTX_FILE), f"{base_name}.gen.pptx")

    for i, (slide, slide_text) in enumerate(zip(presentation.slides, slides), start=1):
        print(f"\nProcessing Slide {i}:")
        print(slide_text)

        description = description_prompt_text(slide_text)
        if not description:
            print(f"Skipping slide {i} due to missing description.")
            continue

        print(f"Generated Description: {description}")

        original_image_path = generate_image(description)
        if original_image_path:
            print(f"Image for Slide {i} saved at: {original_image_path}")

            copy_image_to_clipboard(original_image_path)

            print("Running Siri shortcut to remove background.\n!!!Run it manually the first time!!!")
            run_siri_shortcut(REMOVE_BACKGROUND_SHORTCUT)

            time.sleep(3)
            modified_image_path = get_latest_downloaded_image()

            if modified_image_path:
                print(f"Background-removed image found: {modified_image_path}")
            else:
                print("No background-removed image found, using the original.")

            if addOriginal:
                insert_image_to_slide(slide, original_image_path, Cm(16), Cm(0))
            
            if addSeparated and modified_image_path:
                insert_image_to_slide(slide, modified_image_path, Cm(16), Cm(0))  # Position adjusted for second image

            presentation.save(output_pptx_file)
            print(f"\nPartly saved as: {output_pptx_file}")
        else:
            print(f"Failed to generate image for Slide {i}.")
            
    # Save updated PowerPoint file
    presentation.save(output_pptx_file)
    print(f"\nFinished Generating, saved as: {output_pptx_file}")